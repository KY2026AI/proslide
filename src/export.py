#!/usr/bin/env python3
"""
ProSlide Export Script
HTML Preview → HD Screenshot → PPTX

Requirements:
    pip install playwright python-pptx
    playwright install chromium

Usage:
    python export.py --input preview.html --output slide.pptx
"""

import argparse
import os
import sys
from pathlib import Path

try:
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print(f"Missing dependency: {e}")
    print("Please run: pip install playwright python-pptx && playwright install chromium")
    sys.exit(1)


def export_html_to_pptx(html_path: str, output_path: str, selector: str = ".slide") -> str:
    """
    Export an HTML slide to a PPTX file.

    Args:
        html_path: Path to the HTML preview file.
        output_path: Path to save the generated PPTX.
        selector: CSS selector for the slide element (default: .slide).

    Returns:
        Absolute path to the saved PPTX file.
    """
    html_path = os.path.abspath(html_path)
    if not os.path.exists(html_path):
        raise FileNotFoundError(f"HTML file not found: {html_path}")

    output_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(output_dir, exist_ok=True)

    screenshot_path = os.path.join(output_dir, "slide_export_temp.png")

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        page.goto(f"file://{html_path}")
        page.wait_for_timeout(1000)

        slide = page.query_selector(selector)
        if not slide:
            browser.close()
            raise RuntimeError(f"Element '{selector}' not found in {html_path}")

        slide.screenshot(path=screenshot_path, type="png", scale="device")
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank_layout)
    s.shapes.add_picture(screenshot_path, Inches(0), Inches(0), width=Inches(13.333))
    prs.save(output_path)

    if os.path.exists(screenshot_path):
        os.remove(screenshot_path)

    return os.path.abspath(output_path)


def main():
    parser = argparse.ArgumentParser(description="ProSlide: HTML → Screenshot → PPTX")
    parser.add_argument("--input", "-i", required=True, help="Path to HTML preview file")
    parser.add_argument("--output", "-o", required=True, help="Path to output PPTX file")
    parser.add_argument("--selector", "-s", default=".slide", help="CSS selector for slide element")
    args = parser.parse_args()

    try:
        result = export_html_to_pptx(args.input, args.output, args.selector)
        print(f"PPTX exported successfully: {result}")
    except Exception as e:
        print(f"Export failed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
